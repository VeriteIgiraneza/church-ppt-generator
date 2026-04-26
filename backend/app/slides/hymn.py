"""Hymn slides — one title slide + one slide per verse."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from app.models.hymn import Hymn
from app.slides.base import (
    ALIGN_CENTER,
    ANCHOR_MIDDLE,
    BLACK,
    BLANK_LAYOUT_INDEX,
    FONT_NAME,
    HYMN_UNDERLINE,
    add_centered_paragraph,
    hymn_verse_font_size,
)


def add_hymn_slides(prs: Presentation, hymn: Hymn) -> None:
    """Add all slides for one hymn: title slide + a slide per remaining verse."""
    verses = [v.strip() for v in hymn.verses.split("||") if v.strip()]
    if not verses:
        return

    # First slide: title + verse 1
    _add_title_slide(prs, hymn, verses[0].replace("|", "\n"))

    # Remaining verses: one slide each
    for i in range(1, len(verses)):
        is_last = i == len(verses) - 1
        _add_verse_slide(prs, verses[i].replace("|", "\n"), is_last_verse=is_last)


def _add_title_slide(prs: Presentation, hymn: Hymn, first_verse: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    title_text = f"{hymn.hymn_id} {hymn.title}" if hymn.hymn_id else hymn.title

    # Title
    title_box = slide.shapes.add_textbox(Inches(2), Inches(0.2), Inches(9.33), Inches(1.5))
    add_centered_paragraph(title_box.text_frame, title_text, Pt(40), bold=True)

    # Underline separator
    line_box = slide.shapes.add_textbox(Inches(2), Inches(0.55), Inches(9.33), Inches(1.5))
    add_centered_paragraph(line_box.text_frame, HYMN_UNDERLINE, Pt(30))

    # First verse
    content_box = slide.shapes.add_textbox(Inches(2), Inches(1.2), Inches(9.33), Inches(3.5))
    text_frame = content_box.text_frame
    text_frame.vertical_anchor = ANCHOR_MIDDLE
    text_frame.word_wrap = True
    add_centered_paragraph(text_frame, first_verse, hymn_verse_font_size(len(first_verse)))


def _add_verse_slide(
    prs: Presentation, verse_text: str, *, is_last_verse: bool
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    # Underline at top
    line_box = slide.shapes.add_textbox(Inches(2), Inches(0.56), Inches(9.33), Inches(0.5))
    add_centered_paragraph(line_box.text_frame, HYMN_UNDERLINE, Pt(30))

    # Verse content
    content_box = slide.shapes.add_textbox(Inches(2), Inches(1.2), Inches(9.33), Inches(3.5))
    text_frame = content_box.text_frame
    text_frame.vertical_anchor = ANCHOR_MIDDLE
    text_frame.word_wrap = True
    add_centered_paragraph(text_frame, verse_text, hymn_verse_font_size(len(verse_text)))

    # Asterisks on last verse
    if is_last_verse:
        ast_box = slide.shapes.add_textbox(Inches(6), Inches(6), Inches(9.33), Inches(1))
        ast_frame = ast_box.text_frame
        ast_frame.margin_left = 0
        ast_frame.margin_right = 0
        p = ast_frame.add_paragraph()
        p.text = "*****"
        p.alignment = ALIGN_CENTER
        p.font.size = Pt(30)
        p.font.name = FONT_NAME
        p.font.color.rgb = BLACK