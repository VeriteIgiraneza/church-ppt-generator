from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from app.slides.base import (
    ANCHOR_MIDDLE,
    BLANK_LAYOUT_INDEX,
    add_centered_paragraph,
)


def add_title_slide(
    prs: Presentation, title: str, bible_reference: str = ""
) -> None:
    """Add a title slide with the service title and (optionally) a Bible reference."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    # Service title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(5), Inches(13.13), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.vertical_anchor = ANCHOR_MIDDLE
    add_centered_paragraph(title_frame, title, Pt(45), bold=True)

    # Bible reference (subtitle)
    if bible_reference:
        ref_box = slide.shapes.add_textbox(Inches(0.4), Inches(6), Inches(13.13), Inches(1))
        ref_frame = ref_box.text_frame
        ref_frame.vertical_anchor = ANCHOR_MIDDLE
        add_centered_paragraph(ref_frame, bible_reference, Pt(40))