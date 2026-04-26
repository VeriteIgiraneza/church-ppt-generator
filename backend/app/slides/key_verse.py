"""Key verse slide — highlighted single passage with reference header."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from app.models.bible import BibleVerse
from app.slides.base import (
    ANCHOR_MIDDLE,
    BLANK_LAYOUT_INDEX,
    add_centered_paragraph,
)


def add_key_verse_slide(
    prs: Presentation, verses: list[BibleVerse], reference: str
) -> None:
    """Add a single 'KEY VERSE' slide with the verses joined into one quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    # Header: "KEY VERSE: <reference>"
    title_box = slide.shapes.add_textbox(Inches(0.1), Inches(0.5), Inches(13.13), Inches(1))
    add_centered_paragraph(
        title_box.text_frame, f"KEY VERSE: {reference}", Pt(40), bold=True
    )

    # Quoted verse text
    joined = " ".join(f"{v.verse} {v.text}" for v in verses)

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.33), Inches(3))
    text_frame = content_box.text_frame
    text_frame.vertical_anchor = ANCHOR_MIDDLE
    text_frame.word_wrap = True
    add_centered_paragraph(text_frame, f'"{joined}"', Pt(34))