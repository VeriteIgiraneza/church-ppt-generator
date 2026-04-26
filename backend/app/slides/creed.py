"""Creed slides — title + content, paginated if too long."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from app.models.creed import Creed
from app.slides.base import (
    ALIGN_CENTER,
    ANCHOR_MIDDLE,
    ANCHOR_TOP,
    BLACK,
    BLANK_LAYOUT_INDEX,
    FONT_NAME,
    add_centered_paragraph,
)

# How many "||" paragraph blocks to fit on one slide before paging.
MAX_PARAGRAPHS_PER_SLIDE = 2


def add_creed_slides(prs: Presentation, creed: Creed) -> None:
    """Add slides for a creed, paginating if it doesn't fit on one slide."""
    paragraphs = [p.strip() for p in creed.content.split("||") if p.strip()]

    page: list[str] = []
    for i, paragraph in enumerate(paragraphs):
        page.append(paragraph.replace("|", "\n"))

        is_last = i == len(paragraphs) - 1
        if len(page) >= MAX_PARAGRAPHS_PER_SLIDE or is_last:
            _add_one_creed_slide(prs, creed.name, "\n\n".join(page))
            page = []


def _add_one_creed_slide(prs: Presentation, title: str, content: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    # Title
    title_box = slide.shapes.add_textbox(Inches(3), Inches(0.001), Inches(13.13), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.vertical_anchor = ANCHOR_TOP
    add_centered_paragraph(title_frame, title, Pt(50), bold=True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.01), Inches(1.5), Inches(13.13), Inches(5))
    text_frame = content_box.text_frame
    text_frame.vertical_anchor = ANCHOR_MIDDLE
    text_frame.word_wrap = True

    # Custom paragraph (line_spacing 1.0 differs from add_centered_paragraph default of 1.2)
    p = text_frame.add_paragraph()
    p.text = content.replace("|", "\n")
    p.alignment = ALIGN_CENTER
    p.font.size = Pt(30)
    p.font.name = FONT_NAME
    p.font.color.rgb = BLACK
    p.line_spacing = 1.0