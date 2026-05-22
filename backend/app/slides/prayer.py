"""Prayer slides — title only, except Representative Prayer which adds 'Led by'."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from app.slides.base import (
    ANCHOR_MIDDLE,
    BLANK_LAYOUT_INDEX,
    add_centered_paragraph,
)

REPRESENTATIVE_PRAYER = "Representative Prayer"


def add_prayer_slide(
    prs: Presentation, prayer_type: str, prayer_leader: str | None = None
) -> None:
    """Add a single prayer slide. Shows leader name only for Representative Prayer."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    # Prayer title
    title_box = slide.shapes.add_textbox(Inches(0.1), Inches(2), Inches(13.13), Inches(2))
    title_frame = title_box.text_frame
    title_frame.vertical_anchor = ANCHOR_MIDDLE
    add_centered_paragraph(title_frame, prayer_type, Pt(48), bold=True)

    # Leader (only for Representative Prayer)
    if prayer_leader and prayer_type == REPRESENTATIVE_PRAYER:
        leader_box = slide.shapes.add_textbox(
            Inches(0.1), Inches(4.5), Inches(13.13), Inches(1.5)
        )
        leader_frame = leader_box.text_frame
        leader_frame.vertical_anchor = ANCHOR_MIDDLE
        add_centered_paragraph(leader_frame, f"Led by: {prayer_leader}", Pt(28))