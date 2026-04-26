"""Bible reading slides — verse-number bolding, paginated by character count."""

from pptx.presentation import Presentation
from pptx.util import Inches, Pt

from app.models.bible import BibleVerse
from app.slides.base import (
    ALIGN_LEFT,
    ANCHOR_MIDDLE,
    BLACK,
    BLANK_LAYOUT_INDEX,
    FONT_NAME,
)


def add_bible_reading_slides(
    prs: Presentation,
    verses: list[BibleVerse],
    max_chars_per_slide: int = 900,
) -> None:
    """Add Bible reading slides, paginating verses based on character count."""
    page: list[BibleVerse] = []
    page_chars = 0

    for verse in verses:
        verse_chars = len(f"{verse.verse} {verse.text} ")

        if page_chars + verse_chars > max_chars_per_slide and page:
            _add_one_bible_slide(prs, page)
            page = [verse]
            page_chars = verse_chars
        else:
            page.append(verse)
            page_chars += verse_chars

    if page:
        _add_one_bible_slide(prs, page)


def _add_one_bible_slide(prs: Presentation, verses: list[BibleVerse]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])

    content_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.01), Inches(13.0), Inches(7))
    text_frame = content_box.text_frame
    text_frame.vertical_anchor = ANCHOR_MIDDLE
    text_frame.word_wrap = True

    paragraph = text_frame.add_paragraph()
    paragraph.alignment = ALIGN_LEFT
    paragraph.line_spacing = 1.2

    for verse in verses:
        text = verse.text
        first_chars = text[:3] if len(text) >= 3 else text
        rest = text[3:] if len(text) > 3 else ""

        # Verse number (bold)
        num_run = paragraph.add_run()
        num_run.text = str(verse.verse)
        _style_run(num_run, bold=True)

        # First 3 chars glued to number with non-breaking space (prevents orphans)
        glue_run = paragraph.add_run()
        glue_run.text = f"\u00a0{first_chars}"
        _style_run(glue_run)

        # Remainder
        if rest:
            rest_run = paragraph.add_run()
            rest_run.text = f"{rest} "
            _style_run(rest_run)
        else:
            space_run = paragraph.add_run()
            space_run.text = " "
            _style_run(space_run)


def _style_run(run, *, bold: bool = False) -> None:
    run.font.size = Pt(28)
    run.font.name = FONT_NAME
    run.font.bold = bold
    run.font.color.rgb = BLACK